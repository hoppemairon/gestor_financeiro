import os
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

def formatar_valor_br(valor):
    """Formata um valor numérico para o formato brasileiro (R$)"""
    if pd.isna(valor):
        return ""
    if isinstance(valor, (int, float)):
        return f"R$ {valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    if isinstance(valor, str):
        try:
            valor_num = float(valor.replace(".", "").replace(",", ".").replace("R$", "").replace("R\\$", "").strip())
            return f"R$ {valor_num:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
        except:
            return valor
    return valor

def gerar_apresentacao_vyco(empresa_nome, dados_ppt):
    """
    Gera uma apresentação PPTX em memória com os dados extraídos do Vyco.
    Retorna o objeto io.BytesIO() contendo o arquivo.
    
    dados_ppt esperado:
    {
        "faturamento_estoque": { df_faturamento: DataFrame, df_estoque: DataFrame },
        "projecoes": { df_fluxo: DataFrame }, 
        "parecer_diagnostico": "texto",
        "analise_gpt": "texto", # opcional
        "parecer_antigravity": "texto", # opcional
        "executivo_mensal": DataFrame # opcional
    }
    """
    prs = Presentation()
    
    # Slides Layouts (default template)
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]
    TITLE_ONLY_LAYOUT = prs.slide_layouts[5]
    BLANK_LAYOUT = prs.slide_layouts[6]

    # --- Slide 1: Capa ---
    slide_capa = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide_capa.shapes.title
    subtitle = slide_capa.placeholders[1]
    
    title.text = f"Relatório Financeiro DRE e Fluxo de Caixa"
    subtitle.text = f"Empresa: {empresa_nome}\nGerado em: {datetime.now().strftime('%d/%m/%Y')}"

    # --- Slide 2: Faturamento e Estoque ---
    if "faturamento_estoque" in dados_ppt:
        dados_fe = dados_ppt["faturamento_estoque"]
        if dados_fe.get("df_faturamento") is not None or dados_fe.get("df_estoque") is not None:
            slide_fe = prs.slides.add_slide(TITLE_ONLY_LAYOUT)
            slide_fe.shapes.title.text = "Faturamento e Estoque Mensal"
            
            # Left table - Faturamento
            df_fat = dados_fe.get("df_faturamento")
            if df_fat is not None and not df_fat.empty:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(4.0)
                height = Inches(0.8)
                
                # Filter to only get total columns if we want, or transpose 
                # Converting index back to column for easy table rendering
                df_fat_reset = df_fat.reset_index()
                df_fat_reset.columns = ['Mês', 'Faturamento'] if len(df_fat_reset.columns) <= 2 else df_fat_reset.columns
                
                rows, cols = df_fat_reset.shape
                table = slide_fe.shapes.add_table(min(rows+1, 15), cols, left, top, width, height).table
                
                # Header
                for i, col_name in enumerate(df_fat_reset.columns):
                    table.cell(0, i).text = str(col_name)
                    
                # Body (limit to first 14 months to avoid overlapping)
                for r in range(min(rows, 14)):
                    for c in range(cols):
                        val = df_fat_reset.iloc[r, c]
                        if c > 0 and isinstance(val, (int, float)):
                            table.cell(r+1, c).text = formatar_valor_br(val)
                        else:
                            table.cell(r+1, c).text = str(val)[:30]
            
            # Right table - Estoque 
            df_est = dados_fe.get("df_estoque")
            if df_est is not None and not df_est.empty:
                left = Inches(5.0)
                top = Inches(1.5)
                width = Inches(4.5)
                height = Inches(0.8)
                
                df_est_reset = df_est.reset_index()
                df_est_reset.columns = ['Mês', 'Estoque'] if len(df_est_reset.columns) <= 2 else df_est_reset.columns
                
                rows, cols = df_est_reset.shape
                table_est = slide_fe.shapes.add_table(min(rows+1, 15), cols, left, top, width, height).table
                
                for i, col_name in enumerate(df_est_reset.columns):
                    table_est.cell(0, i).text = str(col_name)
                    
                for r in range(min(rows, 14)):
                    for c in range(cols):
                        val = df_est_reset.iloc[r, c]
                        if c > 0 and isinstance(val, (int, float)):
                            table_est.cell(r+1, c).text = formatar_valor_br(val)
                        else:
                            table_est.cell(r+1, c).text = str(val)[:30]

    # --- Slide 3: Projeções (DRE Simplificado/Resultados) ---
    if "projecoes" in dados_ppt:
        df_proj = dados_ppt["projecoes"].get("df_fluxo")
        if df_proj is not None and not df_proj.empty:
            slide_proj = prs.slides.add_slide(TITLE_ONLY_LAYOUT)
            slide_proj.shapes.title.text = "Resultados - DRE Consolidado"
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9.0)
            height = Inches(5.0)
            
            df_slice = df_proj.copy()
            # Only pick rows that are important for high-level views
            linhas_importantes = ["FATURAMENTO", "RECEITAS", "RECEITA", "IMPOSTOS", "LUCRO OPERACIONAL", "RESULTADO", "ESTOQUE"]
            df_slice = df_slice[df_slice.index.isin(linhas_importantes)]
            
            if not df_slice.empty:
                df_slice.reset_index(inplace=True)
                # Keep first 6 columns if it's too wide
                df_slice = df_slice.iloc[:, :7]
                
                rows, cols = df_slice.shape
                table = slide_proj.shapes.add_table(rows+1, cols, left, top, width, height).table
                
                # Header
                for i, col_name in enumerate(df_slice.columns):
                    table.cell(0, i).text = str(col_name)
                    # format header cells font
                    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(12)
                    
                # Body
                for r in range(rows):
                    for c in range(cols):
                        val = df_slice.iloc[r, c]
                        if c > 0 and (isinstance(val, (int, float)) or (isinstance(val, str) and str(val).replace('.','').replace('-','').isdigit())):
                            try:
                                table.cell(r+1, c).text = formatar_valor_br(float(val))
                            except:
                                table.cell(r+1, c).text = str(val)
                        else:
                            table.cell(r+1, c).text = str(val)[:25]
                        table.cell(r+1, c).text_frame.paragraphs[0].font.size = Pt(11)


    # --- Função auxiliar add texto longo em slides múltiplos ---
    def adicionar_slides_texto(prs, titulo, texto_longo):
        """Quebra um texto longo em múltiplos slides para não vazar a tela"""
        if not texto_longo or pd.isna(texto_longo):
            return
            
        # Simplista: recorta o texto a cada 800 caracteres em blocos
        linhas = str(texto_longo).split('\n')
        chunks = []
        chunk_atual = ""
        for linha in linhas:
            if len(chunk_atual) + len(linha) > 850: # Limit per slide
                chunks.append(chunk_atual)
                chunk_atual = linha + "\n"
            else:
                chunk_atual += linha + "\n"
        
        if chunk_atual.strip():
            chunks.append(chunk_atual)
            
        for i, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            nome_titulo = titulo if i == 0 else f"{titulo} (Cont.)"
            title.text = nome_titulo
            
            # Limpar formatações markdown do texto para o PPT (negrito simples)
            texto_limpo = chunk.replace('**', '').replace('###', '->')
            body.text = texto_limpo
            body.text_frame.paragraphs[0].font.size = Pt(14)
            for p in body.text_frame.paragraphs:
                p.font.size = Pt(14)

    # --- Slide 4: Parecer Diagnóstico ---
    if "parecer_diagnostico" in dados_ppt and dados_ppt["parecer_diagnostico"]:
        adicionar_slides_texto(prs, "Parecer Diagnóstico", dados_ppt["parecer_diagnostico"])

    # --- Slide 5: Análise GPT (Condicional) ---
    if "analise_gpt" in dados_ppt and dados_ppt["analise_gpt"]:
        adicionar_slides_texto(prs, "Análise de IA (GPT)", dados_ppt["analise_gpt"])
        
    # --- Slide 6: Parecer Antigravity (Condicional) ---
    if "parecer_antigravity" in dados_ppt and dados_ppt["parecer_antigravity"]:
        adicionar_slides_texto(prs, "Parecer Antigravity Financeiro", dados_ppt["parecer_antigravity"])

    # --- Slide 7: Relatório Executivo Mensal (Condicional) ---
    if "executivo_mensal" in dados_ppt:
        df_exec = dados_ppt["executivo_mensal"]
        if df_exec is not None and not df_exec.empty:
            slide_exec = prs.slides.add_slide(TITLE_ONLY_LAYOUT)
            slide_exec.shapes.title.text = "Dashboard Executivo - Indicadores Consolidados"
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9.0)
            height = Inches(4.0)
            
            rows, cols = df_exec.shape
            table = slide_exec.shapes.add_table(min(rows+1, 15), cols, left, top, width, height).table
            
            # Header
            for i, col_name in enumerate(df_exec.columns):
                table.cell(0, i).text = str(col_name)
                table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(11)
                
            # Body
            for r in range(min(rows, 14)): # Limite de linhas
                for c in range(cols):
                    val = df_exec.iloc[r, c]
                    table.cell(r+1, c).text = str(val)[:40]
                    table.cell(r+1, c).text_frame.paragraphs[0].font.size = Pt(10)

    # Save to BytesIO stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    
    return ppt_stream
